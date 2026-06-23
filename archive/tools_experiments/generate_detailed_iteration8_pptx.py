#!/usr/bin/env python3
"""
Detailed Iteration 8 Slides for PPT (very explanatory version)

사용자가 요청한 "설명 아주 자세히" 버전.
기존 진화 스토리텔링 스타일에 맞춰:
- 당시 문제
- 우리가 무엇을 했고 왜 했는지
- 어떻게 작동하는지 (메커니즘)
- 실제 결과 + 해석
- 이게 의미하는 것 (패러다임 전환)
을 길고 자세하게 설명하는 슬라이드 4장 생성.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK = RGBColor(0x1E, 0x27, 0x61)
TEAL = RGBColor(0x0F, 0x4C, 0x5C)
ACCENT = RGBColor(0x1A, 0x8A, 0x9B)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
YELLOW = RGBColor(0xF9, 0xA8, 0x25)
RED_DARK = RGBColor(0x5D, 0x2E, 0x2E)
GREEN_DARK = RGBColor(0x1E, 0x4D, 0x3A)

def add_text_box(slide, left, top, width, height, text, font_size=12, bold=False, color=TEXT_DARK, italic=False, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.alignment = align
    return tx

def add_rounded_box(slide, left, top, width, height, fill_color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.fill.background()
    return box

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ========== SLIDE 1: 한계 선언 (아주 자세한 설명) ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    add_text_box(slide1, 0.5, 0.25, 9, 0.4, "Iteration 8 — 새로운 패러다임의 시작", 16, False, RGBColor(0xA0, 0xD2, 0xDB))
    add_text_box(slide1, 0.5, 0.6, 9, 0.7, "왜 'RGB Weight 최적화' 게임은 이제 한계에 도달했는가?", 26, True, WHITE)

    # 문제 상세 설명 박스
    add_rounded_box(slide1, 0.5, 1.4, 9, 3.9, RGBColor(0x16, 0x21, 0x3E))

    add_text_box(slide1, 0.7, 1.55, 8.6, 0.35, "1/2/3 비교 실험에서 드러난 명확한 한계", 14, True, YELLOW)

    detailed_text = """지금까지 우리가 수개월 동안 해온 작업의 본질은 결국 이것이었습니다.

"RGB 세 채널을 어떤 선형 가중치(w_r, w_g, w_b)로 결합하면 강아지 심박 신호가 가장 잘 나오는가?"

우리는 이 질문에 답하기 위해:
• v1 (combined_correct), v2 (60 window), high-HR focused (44 window) 등 여러 버전의 weight를 학습
• Temporal Prior (ramping 0.2→0.55 + re-scan)
• Adaptive weight switching (high/low HR)
• Ensemble (dog_learned + POS + CHROM + ICA 품질 가중)

까지 모두 동원했습니다.

그 결과? 현실적인 Temporal Tracking 모드에서 최고 조합(1+2+3)조차 전체 MAE 36.0에 머물렀고, 특히 Video 3과 Video 7에서는 50~80 bpm 오차가 반복적으로 발생했습니다.

결론: 더 이상 weight를 어떻게 조정하든, peak ambiguity 문제(고심박 cardiac peak vs 100bpm panting artifact)를 근본적으로 해결할 수 없었습니다."""

    add_text_box(slide1, 0.7, 1.95, 8.6, 3.2, detailed_text, 11, False, RGBColor(0xE0, 0xE0, 0xE0))

    # ========== SLIDE 2: Spectrum Selector 상세 메커니즘 ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = LIGHT_BG
    bg2.line.fill.background()

    add_text_box(slide2, 0.5, 0.2, 9, 0.4, "Direction 1: Spectrum-Domain Learned Selector", 18, True, DARK)
    add_text_box(slide2, 0.5, 0.55, 9, 0.35, "RGB Weight를 완전히 버리고, 스펙트럼의 '모양' 자체를 학습하다", 13, True, ACCENT)

    # 핵심 아이디어 박스
    add_rounded_box(slide2, 0.5, 0.95, 9, 1.1, RGBColor(0xE8, 0xF4, 0xF8))
    add_text_box(slide2, 0.65, 1.0, 8.7, 1.0, 
"""기존: RGB (T,3) → [w_r, w_g, w_b] 학습 → 1차원 pulse 신호 → FFT → 가장 큰 peak 선택 (+ Prior penalty)
새로운 접근: RGB (T,3) → 여러 고정 views (Green, G-R, crude POS-like) → 각 view의 주파수 분포(binned spectrum) + shape descriptors (peak 위치, high-HR band vs artifact band power ratio) → 작은 Ridge 모델이 "이 스펙트럼 모양이면 BPM이 얼마일 가능성이 높은가"를 직접 예측""", 10, False, TEXT_DARK)

    # 왜 효과가 있는가 상세 설명
    add_text_box(slide2, 0.5, 2.15, 9, 0.3, "왜 이게 기존 weight 방식보다 강아지 데이터에서 더 잘 맞았는가?", 12, True, TEAL)

    reason_text = """강아지 털(fur)은 인간 피부와 광학적 특성이 완전히 다릅니다. 
인간용으로 설계된 CHROM, POS 같은 linear skin reflection 모델은 강아지에서는 본질적으로 부정확합니다.

더 중요한 점: 고심박(180~210 bpm) 상황에서 cardiac 신호와 panting artifact(~100 bpm harmonic)가 스펙트럼상에서 매우 가깝게 붙어 있습니다.
이때 단순히 "가장 큰 peak"를 고르는 방식은 자주 실패합니다.

Spectrum descriptors는 이 경쟁 상황 자체를 feature로 포착합니다.
"190~210 구간에 에너지가 어느 정도 있고, 95~115 artifact 구간 대비 그 비율이 얼마인가"라는 상대적 shape 정보를 모델이 학습할 수 있게 됩니다.
이것이 60개 데이터만으로도 Video 3에서 8.9 bpm이라는 극적인 결과를 만든 핵심 이유입니다."""

    add_text_box(slide2, 0.5, 2.45, 9, 1.6, reason_text, 10, False, TEXT_DARK)

    # 결과 하이라이트
    add_rounded_box(slide2, 0.5, 4.15, 9, 1.3, RGBColor(0xE8, 0xF8, 0xE8))
    add_text_box(slide2, 0.65, 4.25, 8.7, 0.3, "실제 결과 (7비디오, 동일 sampling)", 11, True, SUCCESS)
    add_text_box(slide2, 0.65, 4.55, 8.7, 0.8, 
"pure1 Spectrum (descriptors only): 전체 MAE 18.2 bpm\n• Video 3 (210 bpm): 8.9 bpm  ← 역대 최고 수준\n• Video 7 (189.5 bpm): 38.9 bpm\n• Video 8 (110.5 bpm): 1.6 bpm\nCV (5-fold, descriptors): 29.5 bpm (60개 window 기준)", 10, False, TEXT_DARK)

    # ========== SLIDE 3: 실제 스펙트럼 이미지 + 해석 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = LIGHT_BG
    bg3.line.fill.background()

    add_text_box(slide3, 0.5, 0.15, 9, 0.35, "실제 데이터로 보는 '왜 Spectrum인가'", 17, True, DARK)

    # 스펙트럼 이미지 임베드
    try:
        slide3.shapes.add_picture("reports/rppg_pet_keypoints/bvp_visualization/3_spectrum.png", Inches(5.3), Inches(0.55), width=Inches(4.4))
    except:
        add_text_box(slide3, 5.3, 1.5, 4.4, 1.5, "[Video 3 Spectrum 이미지]\n고심박 상황에서 cardiac peak와 artifact가 경쟁하는 전형적 모습", 10, False, RGBColor(0x5D, 0x6D, 0x7E), align=PP_ALIGN.CENTER)

    add_text_box(slide3, 0.5, 0.55, 4.7, 0.3, "Video 3 실제 스펙트럼 (target 210 bpm)", 11, True, TEAL)

    interpretation = """이 스펙트럼에서 기존 weight 기반 방법들이 자주 실패했던 이유:

보이는 것처럼 100bpm 근처에 매우 강한 artifact 피크가 있고, 200bpm 근처에 cardiac 신호가 있습니다.
단순 FFT argmax + Prior 페널티만으로는, artifact가 더 강하게 나오는 window에서는 자주 100bpm대에 고착됩니다.

Spectrum 모델은 다릅니다.
"전체 스펙트럼 모양에서 165~225 구간의 상대적 에너지 비율이 얼마인가", "peak 위치가 200 근처인가" 등의 다차원 shape 정보를 동시에 보고 판단합니다.

결과적으로 같은 window에서도 weight 방법이 148 bpm이나 126 bpm으로 잘못 예측하던 상황에서, Spectrum descriptors는 218.9 bpm으로 훨씬 가까운 추정을 해냈습니다."""

    add_text_box(slide3, 0.5, 2.6, 4.7, 2.4, interpretation, 10, False, TEXT_DARK)

    # 하단 핵심 메시지
    add_rounded_box(slide3, 0.5, 5.0, 9, 0.5, DARK)
    add_text_box(slide3, 0.65, 5.05, 8.7, 0.4, "이 한 장의 스펙트럼이 증명하는 것: 강아지 rPPG는 이제 '더 좋은 weight 3개'가 아니라 '스펙트럼의 의미를 읽는 모델'이 필요하다.", 10, True, WHITE, align=PP_ALIGN.CENTER)

    # ========== SLIDE 4: Tracker + 최종 메시지 (매우 자세히) ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = LIGHT_BG
    bg4.line.fill.background()

    add_text_box(slide4, 0.5, 0.15, 9, 0.35, "Direction 3 + 종합: 우리가 얻은 가장 중요한 교훈", 17, True, DARK)

    # Tracker 설명
    add_rounded_box(slide4, 0.5, 0.55, 4.4, 1.6, WHITE)
    add_text_box(slide4, 0.65, 0.6, 4.1, 0.3, "Direction 3: Kalman State Tracker", 11, True, TEAL)
    add_text_box(slide4, 0.65, 0.9, 4.1, 1.15, 
"""기존: 단순 IIR prior (이전 BPM의 65% + 현재 35%)
문제: "갑자기 40 bpm 점프"를 막을 수 없음

새로운 Tracker:
• 상태 = [현재 HR, HR 변화 속도(velocity)]
• velocity는 매 스텝 15%씩 감쇠
• Process noise를 개의 실제 HR 변화 패턴에 맞춰 튜닝
• 큰 점프는 covariance를 통해 자연스럽게 억제""", 9, False, TEXT_DARK)

    # 1+3 결과
    add_rounded_box(slide4, 5.1, 0.55, 4.4, 1.6, WHITE)
    add_text_box(slide4, 5.25, 0.6, 4.1, 0.3, "1+3 통합 결과 및 해석", 11, True, RGBColor(0xC0, 0x39, 0x2B))
    add_text_box(slide4, 5.25, 0.9, 4.1, 1.15, 
"""pure1 (Spectrum): 18.2 MAE
pure3 (Tracker + 약한 obs): 41.3 MAE
1+3 조합: 27.4 MAE

해석: Spectrum이 좋은 측정을 주면 Tracker가 안정화에 도움을 주지만, 아직 Spectrum의 uncertainty를 정확히 모델링하지 못해 과도한 smoothing이 발생.

→ 다음 레버리지 포인트: Spectrum 모델이 자신의 불확실도를 잘 추정하게 만드는 것.""", 9, False, TEXT_DARK)

    # 최종 메시지 박스 (가장 중요)
    add_rounded_box(slide4, 0.5, 2.3, 9, 3.1, DARK)
    add_text_box(slide4, 0.7, 2.45, 8.6, 0.3, "발표에서 가장 강조해야 할 한 문장", 12, True, YELLOW)

    final_quote = '''지금까지 우리는 '더 좋은 RGB 가중치 3개'를 찾기 위해 수많은 실험을 했습니다.
1/2/3 비교를 통해 우리는 그 방향의 한계를 명확히 보았습니다.

그래서 우리는 완전히 다른 길을 택했습니다.
RGB를 1차원으로 압축하는 대신, 주파수 영역에서의 '모양' 자체를 feature로 삼고,
그 모양이 의미하는 BPM을 작은 모델이 직접 학습하게 했습니다.

그 결과, 60개 데이터만으로도 이전 weight 기반 방법이 50~70 bpm 오차를 내던 Video 3에서 8.9 bpm이라는 결과를 얻었습니다.

이것은 단순한 수치 개선이 아닙니다.
강아지 rPPG가 '인간 피부 모델의 변형'에서 '강아지 신호의 본질을 직접 읽는 모델'로 전환되는 첫 걸음입니다.'''

    add_text_box(slide4, 0.7, 2.8, 8.6, 2.4, final_quote, 10, False, WHITE)

    # Save
    output = "reports/PET_RPPG_Iteration8_Detailed_Explanation_Slides.pptx"
    prs.save(output)
    print(f"Generated: {output}")
    print("이 파일은 '설명 아주 자세히' 버전입니다. 메인 PPTX에 붙여넣기 추천.")

if __name__ == "__main__":
    main()