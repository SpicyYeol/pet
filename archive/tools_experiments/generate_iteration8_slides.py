#!/usr/bin/env python3
"""
Generate 3 supplementary slides for Iteration 8 (Spectrum + State Tracker)
to be added to the main evolution story PPTX.

These slides directly address the user's request to add the new 1+3 results
and the "beyond simple RGB weights" conclusion to the presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors (consistent with project history)
DARK = RGBColor(0x1E, 0x27, 0x61)       # Navy
TEAL = RGBColor(0x0F, 0x4C, 0x5C)       # Deep teal
ACCENT = RGBColor(0x1A, 0x8A, 0x9B)     # Bright teal
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
YELLOW = RGBColor(0xF9, 0xA8, 0x25)

def set_cell_text(para, text, font_size=11, bold=False, color=TEXT_DARK):
    para.text = text
    para.font.size = Pt(font_size)
    para.font.bold = bold
    para.font.color.rgb = color
    para.alignment = PP_ALIGN.CENTER

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ========== SLIDE 1: Title + Motivation ==========
    slide_layout = prs.slide_layouts[6]  # Blank
    slide1 = prs.slides.add_slide(slide_layout)

    # Dark background
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    # Title
    tx = slide1.shapes.add_textbox(Inches(0.5), Inches(0.7), Inches(9), Inches(0.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "Iteration 8"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0xA0, 0xD2, 0xDB)

    tx2 = slide1.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.9))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "Spectrum-Domain Selector + State Tracker"
    p2.font.size = Pt(30)
    p2.font.bold = True
    p2.font.color.rgb = WHITE

    tx3 = slide1.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(0.5))
    tf3 = tx3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = "단순 RGB Weight 최적화를 넘어선 첫 번째 진짜 대안"
    p3.font.size = Pt(16)
    p3.font.italic = True
    p3.font.color.rgb = RGBColor(0xA0, 0xD2, 0xDB)

    # Problem box
    box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.9), Inches(9), Inches(2.2))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x16, 0x21, 0x3E)
    box.line.fill.background()

    tx4 = slide1.shapes.add_textbox(Inches(0.7), Inches(3.05), Inches(8.6), Inches(0.35))
    tf4 = tx4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "한계 인식 (1/2/3 실험 결과)"
    p4.font.size = Pt(13)
    p4.font.bold = True
    p4.font.color.rgb = YELLOW

    tx5 = slide1.shapes.add_textbox(Inches(0.7), Inches(3.45), Inches(8.6), Inches(1.5))
    tf5 = tx5.text_frame
    tf5.word_wrap = True
    p5 = tf5.paragraphs[0]
    p5.text = "• RGB weight (v1 / high-HR focused / combined_correct) + Prior + Ensemble까지 동원해도 현실적 모드에서 36 MAE 한계\n• 특히 고심박 Video 3/7에서 peak ambiguity (100bpm artifact와의 경쟁)가 근본적으로 해결되지 않음\n• 결론: '더 좋은 weight 3개'를 찾는 게임 자체가 한계에 도달"
    p5.font.size = Pt(12)
    p5.font.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)

    # Simple visual: Old vs New paradigm (mini diagram at bottom)
    # Old box
    old_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.15), Inches(3.8), Inches(0.38))
    old_box.fill.solid()
    old_box.fill.fore_color.rgb = RGBColor(0x5D, 0x2E, 0x2E)
    old_box.line.fill.background()
    tx_old = slide1.shapes.add_textbox(Inches(0.6), Inches(5.18), Inches(3.6), Inches(0.32))
    tf_old = tx_old.text_frame
    p_old = tf_old.paragraphs[0]
    p_old.text = "❌ Old: RGB → 3 weights → 1D pulse → FFT + Prior"
    p_old.font.size = Pt(9)
    p_old.font.color.rgb = RGBColor(0xFF, 0xCC, 0xCC)

    # Arrow
    arr = slide1.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.4), Inches(5.2), Inches(0.7), Inches(0.28))
    arr.fill.solid()
    arr.fill.fore_color.rgb = YELLOW
    arr.line.fill.background()

    # New box
    new_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(5.15), Inches(4.3), Inches(0.38))
    new_box.fill.solid()
    new_box.fill.fore_color.rgb = RGBColor(0x1E, 0x4D, 0x3A)
    new_box.line.fill.background()
    tx_new = slide1.shapes.add_textbox(Inches(5.3), Inches(5.18), Inches(4.1), Inches(0.32))
    tf_new = tx_new.text_frame
    p_new = tf_new.paragraphs[0]
    p_new.text = "✓ New: Spectrum descriptors + Kalman (no RGB weights)"
    p_new.font.size = Pt(9)
    p_new.font.color.rgb = RGBColor(0xC8, 0xE6, 0xC9)

    # ========== SLIDE 2: Spectrum Results + Real Spectrum Image ==========
    slide2 = prs.slides.add_slide(slide_layout)

    # Light bg
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = LIGHT_BG
    bg2.line.fill.background()

    # Title
    tx = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.45))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "Direction 1: Spectrum-Domain Learned Selector"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK

    tx2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.3))
    tf2 = tx2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "RGB Weight 완전 포기 → 스펙트럼 Shape 직접 모델링"
    p2.font.size = Pt(12)
    p2.font.italic = True
    p2.font.color.rgb = ACCENT

    # Key idea box (smaller)
    box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.95), Inches(9), Inches(0.6))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xF8)
    box.line.fill.background()

    tx3 = slide2.shapes.add_textbox(Inches(0.65), Inches(1.02), Inches(8.7), Inches(0.45))
    tf3 = tx3.text_frame
    tf3.word_wrap = True
    p3 = tf3.paragraphs[0]
    p3.text = "고정 views(Green, G-R 등) → binned spectrum + descriptors (peak 위치 + high-HR band vs 100bpm artifact power ratio) → Ridge"
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_DARK

    # Left: Results table title + table (narrower)
    tx4 = slide2.shapes.add_textbox(Inches(0.5), Inches(1.65), Inches(4.8), Inches(0.28))
    tf4 = tx4.text_frame
    p4 = tf4.paragraphs[0]
    p4.text = "7비디오 결과 (동일 sampling)"
    p4.font.size = Pt(11)
    p4.font.bold = True
    p4.font.color.rgb = TEXT_DARK

    # Compact table
    headers = ["Config", "Overall", "V3(210)", "V7"]
    rows = [
        ["Spectrum (1)", "18.2", "8.9 ★", "38.9"],
        ["Weight+Prior+Ens (이전)", "~36", "50~70+", "50~60+"]
    ]
    table_left = 0.5
    col_widths = [2.2, 0.9, 0.9, 0.8]
    row_height = 0.32
    start_y = 1.95

    # Header
    x = table_left
    for header, w in zip(headers, col_widths):
        shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(start_y), Inches(w), Inches(row_height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = TEAL
        shape.line.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
        tx = slide2.shapes.add_textbox(Inches(x + 0.03), Inches(start_y + 0.04), Inches(w - 0.06), Inches(row_height - 0.08))
        set_cell_text(tx.text_frame.paragraphs[0], header, font_size=9, bold=True, color=WHITE)
        x += w

    # Rows
    for ri, row in enumerate(rows):
        y = start_y + (ri + 1) * row_height
        x = table_left
        bg_color = RGBColor(0xD5, 0xF5, 0xE3) if ri == 0 else (WHITE if ri % 2 == 0 else RGBColor(0xF4, 0xF6, 0xF7))
        for cell, w in zip(row, col_widths):
            shape = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(row_height))
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg_color
            shape.line.color.rgb = RGBColor(0xBD, 0xC3, 0xC7)
            is_best = "18.2" in cell or "8.9" in cell
            tx = slide2.shapes.add_textbox(Inches(x + 0.03), Inches(y + 0.04), Inches(w - 0.06), Inches(row_height - 0.08))
            set_cell_text(tx.text_frame.paragraphs[0], cell, font_size=9, bold=is_best)
            x += w

    # Right side: Real spectrum image + caption
    try:
        img_path = "reports/rppg_pet_keypoints/bvp_visualization/3_spectrum.png"
        pic = slide2.shapes.add_picture(img_path, Inches(5.5), Inches(1.6), width=Inches(4.2))
        # Caption under image
        cap = slide2.shapes.add_textbox(Inches(5.5), Inches(4.35), Inches(4.2), Inches(0.5))
        cf = cap.text_frame
        cf.word_wrap = True
        cp = cf.paragraphs[0]
        cp.text = "Video 3 실제 스펙트럼: 210 bpm cardiac band vs ~100 bpm artifact 경쟁 (이게 weight만으로는 풀리지 않았던 이유)"
        cp.font.size = Pt(9)
        cp.font.color.rgb = RGBColor(0x5D, 0x6D, 0x7E)
    except Exception as e:
        # Fallback text if image missing
        fb = slide2.shapes.add_textbox(Inches(5.5), Inches(2.0), Inches(4.2), Inches(2))
        fbf = fb.text_frame
        fbf.word_wrap = True
        fbp = fbf.paragraphs[0]
        fbp.text = "[Video 3 Spectrum 이미지]\n고심박에서 100bpm artifact와 cardiac peak이 경쟁하는 전형적인 상황"
        fbp.font.size = Pt(10)
        fbp.font.color.rgb = RGBColor(0x5D, 0x6D, 0x7E)

    # Bottom insight
    tx5 = slide2.shapes.add_textbox(Inches(0.5), Inches(4.9), Inches(9), Inches(0.5))
    tf5 = tx5.text_frame
    tf5.word_wrap = True
    p5 = tf5.paragraphs[0]
    p5.text = "의미: RGB weight 없이 18.2 MAE. Video 3에서 8.9 (역대 최고 수준). CV 29.5 bpm (60개 데이터)."
    p5.font.size = Pt(11)
    p5.font.color.rgb = SUCCESS
    p5.font.bold = True

    # ========== SLIDE 3: Tracker + Takeaway ==========
    slide3 = prs.slides.add_slide(slide_layout)

    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = LIGHT_BG
    bg3.line.fill.background()

    tx = slide3.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "Direction 3: State Tracker + 종합 교훈"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Left box - Tracker
    box1 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.85), Inches(4.3), Inches(1.7))
    box1.fill.solid()
    box1.fill.fore_color.rgb = WHITE
    box1.line.color.rgb = RGBColor(0xAE, 0xD6, 0xF1)

    tx = slide3.shapes.add_textbox(Inches(0.65), Inches(0.95), Inches(4), Inches(0.3))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "Direction 3: Kalman Tracker"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEAL

    tx2 = slide3.shapes.add_textbox(Inches(0.65), Inches(1.25), Inches(4), Inches(1.2))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "• 2-state (HR + velocity)\n• 개 생리학 제약 명시적 모델링\n• 기존 IIR Prior 완전 대체\n• pure3 (약한 obs): 41.3 MAE\n→ 관측 품질이 핵심"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_DARK

    # Right box - 1+3
    box2 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(0.85), Inches(4.5), Inches(1.7))
    box2.fill.solid()
    box2.fill.fore_color.rgb = WHITE
    box2.line.color.rgb = RGBColor(0xF5, 0xB7, 0xB1)

    tx = slide3.shapes.add_textbox(Inches(5.15), Inches(0.95), Inches(4.2), Inches(0.3))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "1+3 통합 결과"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)

    tx2 = slide3.shapes.add_textbox(Inches(5.15), Inches(1.25), Inches(4.2), Inches(1.2))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "pure1 (Spectrum): 18.2\n1+3: 27.4 (아직 calibration 부족)\n\n→ Spectrum uncertainty를 더 정확히\n   추정하면 1+3이 진짜 강력해질 것"
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_DARK

    # Big takeaway box
    box3 = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.7), Inches(9), Inches(2.4))
    box3.fill.solid()
    box3.fill.fore_color.rgb = DARK
    box3.line.fill.background()

    tx = slide3.shapes.add_textbox(Inches(0.7), Inches(2.85), Inches(8.6), Inches(0.3))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = "핵심 메시지"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = YELLOW

    tx2 = slide3.shapes.add_textbox(Inches(0.7), Inches(3.2), Inches(8.6), Inches(1.7))
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = '"우리는 이제 \'더 좋은 weight 3개\'를 찾는 단계를 넘어,\n강아지 심박 신호의 본질적 특성(스펙트럼 shape + 생리학적 시간 변화)을 직접 모델링하기 시작했습니다."\n\n다음 과제: Spectrum 학습 데이터 대량 확보 (200+ windows) + tiny 1D conv on periodogram + uncertainty calibration'
    p2.font.size = Pt(12)
    p2.font.color.rgb = WHITE

    # Save
    output_path = "reports/PET_RPPG_Iteration8_Spectrum_Tracker_Slides.pptx"
    prs.save(output_path)
    print(f"Generated: {output_path}")
    print("→ 이 파일을 열어서 메인 PPTX(PET_RPPG_Evolution_Story_16slides.pptx)에 슬라이드를 복사하세요.")

if __name__ == "__main__":
    main()