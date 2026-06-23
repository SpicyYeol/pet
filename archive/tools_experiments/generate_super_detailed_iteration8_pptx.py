#!/usr/bin/env python3
"""
Super Detailed Iteration 8 Slides for Non-Experts (비전문가 친화 + 매우 상세 설명 버전)

사용자 요청:
- 더 길게
- 알고리즘 상세히 설명
- 비전문가가 이해할 수 있게 (비유, 단계별 설명, 쉬운 말, 반복 강조)

슬라이드 6장으로 구성:
1. 문제의 본질 (왜 기존 방법이 한계인지, 쉬운 비유로)
2. 새로운 아이디어 소개 (Spectrum Selector가 뭔지, 비유로)
3. 어떻게 작동하는가? (단계별 상세 설명 + 그림)
4. 실제 스펙트럼 이미지로 보는 차이 (Video 3 이미지 + 해설)
5. Kalman Tracker는 왜 필요한가? (시간 흐름 안정화, 비유)
6. 전체 요약 + 우리가 얻은 가장 큰 교훈 (길고 강한 메시지)

모든 설명은 일반인도 이해할 수 있도록 작성.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK = RGBColor(0x1E, 0x27, 0x61)
TEAL = RGBColor(0x0F, 0x4C, 0x5C)
ACCENT = RGBColor(0x1A, 0x8A, 0x9B)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
YELLOW = RGBColor(0xF9, 0xA8, 0x25)
SOFT_BLUE = RGBColor(0xE8, 0xF4, 0xF8)
SOFT_GREEN = RGBColor(0xE8, 0xF8, 0xE8)

def add_text_box(slide, left, top, width, height, text, font_size=11, bold=False, color=TEXT_DARK, italic=False):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    return tx

def add_box(slide, left, top, width, height, fill_color):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    box.line.fill.background()
    return box

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ========== SLIDE 1: 문제의 본질 (비전문가용 긴 설명) ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    add_text_box(slide1, 0.5, 0.2, 9, 0.35, "Iteration 8 — 우리가 왜 완전히 새로운 길을 가기로 했나", 15, False, RGBColor(0xA0, 0xD2, 0xDB))
    add_text_box(slide1, 0.5, 0.55, 9, 0.6, "강아지 심박을 재는 일이 왜 이렇게 어려웠을까?\n(비전문가를 위한 아주 쉬운 설명)", 22, True, WHITE)

    add_box(slide1, 0.5, 1.25, 9, 4.1, RGBColor(0x16, 0x21, 0x3E))

    story = """[비유로 시작하는 문제 설명]

우리가 지금까지 하던 일은, 마치 "강아지가 뛰어다니는 방 안에서 심장 소리를 녹음하는 것"과 비슷했습니다.

문제는 방 안에 '심장 소리' 말고도 '헐떡이는 소리(헐떡임)', '털이 스치는 소리', '조명 깜빡임' 등 온갖 잡음이 가득하다는 점입니다.

기존 방법들은 이렇게 했습니다:
"RGB라는 3개의 마이크로 소리를 녹음해서, 어떤 비율로 섞으면 심장 소리가 가장 크게 들릴까?" 하고 수학적으로 계산해서 섞는 비율(가중치)을 찾는 것이었습니다.

우리는 이 방법을 수없이 반복했습니다.
- 가중치를 여러 번 바꿔서 실험
- 이전 순간의 심박수를 참고해서 이번 순간을 보정 (Prior)
- 여러 방법으로 동시에 재서 투표 (Ensemble)

그런데 Video 3과 Video 7처럼 심장이 정말 빠르게 뛰는 상황(210번, 189번)에서는 이상한 일이 계속 일어났습니다.

헐떡이는 소리(대략 100번 정도)가 심장 소리보다 더 크게 들리는 순간이 많았고, 컴퓨터는 "아, 이게 심장 소리구나!" 하고 잘못 판단하는 경우가 자주 발생했습니다.

결국 "더 좋은 섞는 비율(가중치)"을 아무리 찾아도, 근본적인 문제가 해결되지 않았습니다.

그래서 우리는 이렇게 물었습니다.
"섞는 비율을 찾는 것 자체가 잘못된 접근은 아닐까? 다른 방식으로 접근할 수는 없을까?"

그 답이 바로 이번 Iteration 8입니다."""

    add_text_box(slide1, 0.7, 1.4, 8.6, 3.8, story, 10, False, RGBColor(0xE8, 0xE8, 0xE8))

    # ========== SLIDE 2: Spectrum Selector 소개 (비유 + 원리) ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = LIGHT_BG
    bg2.line.fill.background()

    add_text_box(slide2, 0.5, 0.15, 9, 0.35, "새로운 접근: 스펙트럼 모양으로 보는 방법", 17, True, DARK)
    add_text_box(slide2, 0.5, 0.5, 9, 0.4, "심장 소리를 '크기'가 아니라 '리듬 패턴'으로 판단하기", 13, True, ACCENT)

    add_box(slide2, 0.5, 0.95, 9, 1.5, SOFT_BLUE)
    analogy = """[쉬운 비유]

기존 방법 = "방 안에 있는 여러 마이크 소리를 섞어서, 가장 큰 소리가 나는 순간을 심장 박동이라고 부르자"

새로운 방법 = "심장 소리와 헐떡이는 소리는 리듬이 다르다. 심장은 규칙적으로 빠르게 두근거리고, 헐떡임은 조금 더 느리고 불규칙하다. 그러니 전체 리듬 패턴(스펙트럼)을 보고 '이 패턴은 심장 리듬에 가깝다'고 판단하자"

마치 노래를 들을 때 "이 부분이 가장 시끄럽네"가 아니라 "이 멜로디의 전체 흐름이 발라드인지, 댄스곡인지"를 구분하는 것과 비슷합니다."""
    add_text_box(slide2, 0.65, 1.05, 8.7, 1.3, analogy, 10, False, TEXT_DARK)

    add_text_box(slide2, 0.5, 2.55, 9, 0.3, "실제로 어떻게 다를까? (단계별 비교)", 12, True, TEAL)

    comparison = """기존 Weight 방식 (지금까지 우리가 하던 일)
1. RGB 영상에서 작은 영역(ROI)을 골라서 빨강, 초록, 파랑 신호를 뽑는다
2. "이 신호를 어떻게 섞으면 심장 소리가 잘 들릴까?" 하고 수학으로 가중치 3개를 찾는다
3. 섞어서 하나의 신호를 만든다
4. 이 신호를 주파수로 바꿔서 가장 큰 봉우리를 찾는다
5. 이전 심박수와 비교해서 너무 다르면 보정한다

Spectrum 방식 (이번에 새로 만든 방법)
1. RGB 영상에서 작은 영역을 고르는 것은 같다
2. 하지만 RGB를 섞는 '비율'을 찾지 않는다. 대신 여러 가지 고정된 방식으로 신호를 본다 (초록만 보기, 초록-빨강 차이 보기 등)
3. 각 신호를 주파수로 바꾼다
4. 주파수 그래프 전체의 '모양'을 본다
   - 180~220 사이에 에너지가 얼마나 모여 있는가?
   - 95~115 사이(헐떡임 구간)와 비교해서 어느 쪽이 더 강한가?
5. 이 모양 정보를 작은 수학 모델에 넣어서 "이 모양이면 심박수가 대략 얼마일 가능성이 높다"고 직접 예측한다"""
    add_text_box(slide2, 0.5, 2.85, 9, 2.5, comparison, 9.5, False, TEXT_DARK)

    # ========== SLIDE 3: 왜 이게 강아지에게 더 잘 맞나 ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = LIGHT_BG
    bg3.line.fill.background()

    add_text_box(slide3, 0.5, 0.15, 9, 0.35, "왜 Spectrum 방식이 강아지 데이터에서 더 잘 맞았을까?", 16, True, DARK)

    reason = """[강아지 털이라는 특수한 환경]

강아지는 피부가 아니라 '털'로 덮여 있습니다. 털은 빛을 산란시키고, 색이 있고, 움직일 때마다 신호가 흔들립니다.
인간용으로 만든 CHROM, POS 같은 방법들은 '인간 피부는 이렇게 빛을 반사한다'는 가정에 기반합니다. 강아지에게는 이 가정이 잘 맞지 않습니다.

[고심박 + 헐떡임이라는 치명적인 조합]

강아지가 흥분하거나 운동하면 심박수가 180~220까지 올라갑니다.
그런데 강아지는 동시에 헐떡입니다. 헐떡임은 대략 90~120회 정도의 주기를 만듭니다.
스펙트럼상에서 이 두 주파수가 가까이 붙어있으면, 단순히 "가장 큰 소리"를 찾는 방법은 자주 헷갈립니다.

Spectrum descriptors는 이 헷갈림 자체를 데이터로 활용합니다.
"이번 그래프에서 200 근처가 100 근처보다 얼마나 더 두드러지는가?"라는 상대적인 정보를 학습합니다.
그래서 같은 상황에서 기존 방법이 126 bpm으로 잘못 예측할 때, Spectrum은 210에 훨씬 가까운 답을 내놓을 수 있습니다."""
    add_text_box(slide3, 0.5, 0.55, 9, 2.3, reason, 10, False, TEXT_DARK)

    add_box(slide3, 0.5, 2.95, 9, 2.4, SOFT_GREEN)
    add_text_box(slide3, 0.65, 3.05, 8.7, 0.3, "실제 숫자로 본 차이 (Video 3 예시)", 11, True, SUCCESS)
    add_text_box(slide3, 0.65, 3.35, 8.7, 1.9, 
"""같은 영상 구간에서:

기존 최고 성능 방법 (dog_learned + Prior):
→ 126 bpm 또는 148 bpm 정도로 예측 (실제 210 bpm)

Spectrum descriptors 방법:
→ 218.9 bpm (실제 210 bpm에 매우 가까움)

이 차이가 바로 "크기만 보는 것"과 "전체 모양을 보는 것"의 차이입니다.
비전문가 여러분이 기억하실 한 문장:
"심장 소리가 크다고 해서 항상 심장 소리가 아닐 수 있다. 리듬의 전체 그림을 봐야 진짜를 알 수 있다." """, 10, False, TEXT_DARK)

    # ========== SLIDE 4: 실제 스펙트럼 이미지 설명 ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = LIGHT_BG
    bg4.line.fill.background()

    add_text_box(slide4, 0.5, 0.1, 9, 0.3, "실제 화면으로 보는 차이 — Video 3 스펙트럼", 15, True, DARK)

    try:
        slide4.shapes.add_picture("reports/rppg_pet_keypoints/bvp_visualization/3_spectrum.png", Inches(5.0), Inches(0.45), width=Inches(4.7))
    except:
        pass

    explanation = """이 그래프는 Video 3의 한 구간을 주파수로 바꾼 것입니다.

가로축 = 주파수 (왼쪽이 느린 리듬, 오른쪽이 빠른 리듬)
세로축 = 그 리듬이 얼마나 강한가

빨간색으로 표시한 부분이 실제 강아지 심박(210회)에 해당하는 구간입니다.
파란색으로 표시한 부분이 헐떡임(약 100회) 구간입니다.

기존 방법들은 세로로 가장 높은 막대기를 찾습니다.
이 그래프에서는 파란색 쪽이 더 높게 나오는 경우가 많아서, 컴퓨터가 "아, 이게 심장이다!" 하고 잘못 골라버립니다.

Spectrum 방법은 다릅니다.
그래프 전체를 보고 "빨간색 구간이 파란색 구간보다 상대적으로 얼마나 더 두드러지는가?"를 계산합니다.
그리고 과거에 비슷한 모양을 보였던 데이터들을 기억하고 있다가 "이런 모양일 때는 대개 심박수가 200 근처였다"고 판단합니다.

이 작은 차이가 210 bpm을 126 bpm으로 착각하던 문제를 8.9 bpm 오차로 줄여준 핵심입니다."""
    add_text_box(slide4, 0.5, 0.45, 4.4, 4.8, explanation, 9.5, False, TEXT_DARK)

    # ========== SLIDE 5: Kalman Tracker (비전문가용) ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = LIGHT_BG
    bg5.line.fill.background()

    add_text_box(slide5, 0.5, 0.15, 9, 0.35, "시간의 흐름을 고려하는 Kalman Tracker", 16, True, DARK)
    add_text_box(slide5, 0.5, 0.5, 9, 0.35, "한 순간만 보지 말고, '이전 순간부터 지금까지'의 변화를 보자", 12, True, ACCENT)

    tracker_explain = """[비유]

자동차 내비게이션이 있다고 생각해보세요.
갑자기 "지금 속도가 50km/h인데, 3초 후에 280km/h입니다"라고 말하면 이상하죠.
자동차는 그렇게 급격히 속도를 올릴 수 없습니다.

강아지의 심박수도 마찬가지입니다.
3초 전에 140이었는데 지금 210이라고 갑자기 뛰는 일은 거의 없습니다.
(물론 놀라거나 하면 올라가지만, 그 상승에도 어느 정도 시간이 걸립니다)

기존 Prior는 이 사실을 대략적으로만 반영했습니다.
"이전 값의 65% + 이번 값의 35%를 섞자" 정도였습니다.

Kalman Tracker는 좀 더 똑똑합니다.
- 현재 심박수
- 심박수가 올라가고 있는지, 내려가고 있는지 (변화 속도)
이 두 가지를 동시에 추적합니다.

그리고 "변화 속도는 시간이 지나면 점점 0에 가까워진다(안정된다)"는 규칙을 넣습니다.
그래서 한 번 잘못된 값이 들어와도, 그 다음 순간부터는 점점 현실적인 값으로 돌아오게 됩니다.

이게 바로 '생리학적으로 그럴듯한' 추정을 가능하게 하는 핵심입니다."""
    add_text_box(slide5, 0.5, 0.9, 9, 3.5, tracker_explain, 10, False, TEXT_DARK)

    add_box(slide5, 0.5, 4.5, 9, 0.95, SOFT_GREEN)
    add_text_box(slide5, 0.65, 4.55, 8.7, 0.85, 
"""중요한 점: Tracker는 혼자서 잘하지 못합니다.
Tracker에게 "이번 측정값이 126 bpm입니다"라고 주면, Tracker는 그걸 부드럽게 만들어줄 뿐입니다.
진짜 좋은 측정값을 주는 것이 Spectrum Selector의 역할입니다.
그래서 1(Spectrum) + 3(Tracker)을 함께 쓰는 것이 최종 목표입니다.""", 10, False, TEXT_DARK)

    # ========== SLIDE 6: 최종 메시지 (길고 강하게) ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = DARK
    bg6.line.fill.background()

    add_text_box(slide6, 0.5, 0.2, 9, 0.35, "우리가 얻은 가장 중요한 교훈", 15, False, RGBColor(0xA0, 0xD2, 0xDB))
    add_text_box(slide6, 0.5, 0.55, 9, 0.6, "강아지 rPPG는 이제 새로운 장으로 들어섰습니다", 22, True, WHITE)

    final = """지금까지 우리는 "RGB를 어떻게 섞으면 좋을까?"라는 질문에 수개월을 바쳤습니다.
그 질문 자체가 잘못된 것은 아니었습니다. 그 질문으로 상당한 진전을 이루기도 했습니다.

하지만 1/2/3 실험을 통해 우리는 깨달았습니다.
그 질문의 답을 아무리 열심히 찾아도, 강아지의 실제 심박수를 안정적으로 재는 데는 한계가 있다는 것을.

그래서 우리는 질문을 바꿨습니다.
"강아지의 심박 신호는 어떤 '모양'을 가지고 있는가? 그 모양을 어떻게 읽어야 진짜 심박수를 알 수 있는가?"

이 새로운 질문에 답하기 위해 우리는:
- RGB를 섞는 대신, 주파수 그래프 전체의 모양을 보기로 했습니다.
- 한 순간의 모양뿐만 아니라, 시간에 따라 그 모양이 어떻게 변하는지도 함께 보기로 했습니다.

그 결과, 60개 정도의 데이터만으로도 이전에 50~70 bpm 오차가 나던 구간에서 8~9 bpm 수준까지 오차를 줄일 수 있었습니다.

이 숫자 자체도 중요하지만, 더 중요한 것은 방향입니다.
우리는 이제 "인간을 위해 만든 방법을 강아지에게 억지로 맞추는" 단계에서 벗어나,
"강아지 신호가 가진 고유한 특징을 직접 배우고 이해하는" 단계로 한 발 내딛었습니다.

이 길이 아직 완성된 길은 아닙니다. 데이터도 더 필요하고, 모델도 더 정교해져야 합니다.
하지만 적어도 우리는 올바른 방향으로 첫걸음을 뗐습니다.

강아지의 심장은 우리와 다릅니다.
그 다름을 인정하고, 그 다름에 맞는 방식으로 바라보는 것.
그것이 우리가 이번에 배운 가장 큰 교훈입니다."""

    add_box(slide6, 0.5, 1.25, 9, 4.1, RGBColor(0x16, 0x21, 0x3E))
    add_text_box(slide6, 0.7, 1.4, 8.6, 3.8, final, 10, False, RGBColor(0xE8, 0xE8, 0xE8))

    output = "reports/PET_RPPG_Iteration8_SuperDetailed_NonExpert_Slides.pptx"
    prs.save(output)
    print(f"Generated: {output}")
    print("비전문가도 이해할 수 있도록 아주 길고 자세한 설명으로 구성된 6장 슬라이드입니다.")

if __name__ == "__main__":
    main()