#!/usr/bin/env python3
"""
Visual + Detailed Algorithm Explanation Slides for Iteration 8

사용자 요청: "알고리즘을 그림을 넣어가며 디테일하게 해"

- 각 알고리즘 단계마다 그림(이미지 + shape 다이어그램)과 함께 매우 상세한 설명
- 비전문가도 따라갈 수 있게 단계별로 시각화
- 실제 프로젝트 이미지 다수 임베드 (스펙트럼, 플로우차트)
- python-pptx shapes로 Old vs New 파이프라인, Spectrum Feature, Kalman State 다이어그램 직접 제작

슬라이드 6장 구성 (시각 중심 + 상세 설명):

1. 전체 맥락: 기존 전체 파이프라인 + 우리가 직면한 문제 (플로우차트 이미지 + 강조)
2. Traditional Weight Pipeline 상세 시각화 (block diagram + 문제점)
3. Spectrum Selector Pipeline 상세 시각화 (block diagram + 실제 스펙트럼 이미지)
4. Spectrum Descriptors가 무엇인가? (개념도 + Video 3 스펙트럼 예시 상세 해설)
5. Kalman State Tracker 시각화 (상태 전이 다이어그램 + 비유)
6. Old vs New 완전 비교 + 최종 인사이트 (시각적 요약 + 긴 설명)

모든 설명은 디테일하고, 그림과 텍스트가 함께 배치.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# Colors
DARK = RGBColor(0x1E, 0x27, 0x61)
TEAL = RGBColor(0x0F, 0x4C, 0x5C)
ACCENT = RGBColor(0x1A, 0x8A, 0x9B)
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
SUCCESS = RGBColor(0x27, 0xAE, 0x60)
YELLOW = RGBColor(0xF9, 0xA8, 0x25)
SOFT_BLUE = RGBColor(0xE8, 0xF4, 0xF8)
SOFT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
RED_LIGHT = RGBColor(0xFF, 0xEB, 0xEE)
GREEN_LIGHT = RGBColor(0xE8, 0xF5, 0xE9)

def add_text(slide, left, top, w, h, text, size=10, bold=False, color=TEXT_DARK, italic=False, align=None):
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return tx

def add_box(slide, left, top, w, h, fill):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.fill.background()
    return box

def add_arrow(slide, left, top, w, h, color):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(w), Inches(h))
    arr.fill.solid()
    arr.fill.fore_color.rgb = color
    arr.line.fill.background()
    return arr

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    # ========== SLIDE 1: 전체 맥락 + 플로우차트 ==========
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_BG
    bg.line.fill.background()

    add_text(slide1, 0.5, 0.15, 9, 0.35, "Iteration 8: 알고리즘 전체 흐름과 우리가 직면한 문제", 16, True, DARK)

    # Embed the best flowchart
    try:
        slide1.shapes.add_picture("reports/rppg_pet_keypoints/PET_RPPG_Horizontal_Flowchart_Final.png", 
                                  Inches(0.3), Inches(0.55), width=Inches(9.4))
    except Exception as e:
        add_text(slide1, 0.5, 0.6, 9, 1.5, "[전체 플로우차트 이미지]\nPET_RPPG_Horizontal_Flowchart_Final.png", 10, False, RGBColor(0x5D, 0x6D, 0x7E))

    add_box(slide1, 0.5, 3.3, 9, 2.1, RGBColor(0x16, 0x21, 0x3E))
    add_text(slide1, 0.65, 3.4, 8.7, 0.3, "이 전체 파이프라인에서 가장 어려웠던 부분 (빨간색 강조)", 11, True, YELLOW)
    
    problem = """위 플로우차트에서 "rPPG Signal Extraction"과 "BPM Estimation" 단계가 우리가 수개월 동안 가장 고생한 부분입니다.

기존에는 이 두 단계를 "RGB를 어떻게 섞을까?(Weight)" + "가장 큰 피크를 어떻게 고를까?(Prior)"로 해결하려 했습니다.
하지만 고심박(Video 3, 7)에서 헐떡임 artifact(100bpm 근처)와 cardiac 신호(180~210bpm)가 스펙트럼상에서 너무 가까워서, 
아무리 weight를 잘 찾아도 "잘못된 피크"를 강하게 선택하는 경우가 계속 발생했습니다.

결과: 1/2/3 실험에서 현실적 조건으로 최고 36 MAE. 특히 Video 3/7에서 50~80+ 오차 반복.

→ 그래서 우리는 "RGB를 섞는 방식 자체"를 바꾸기로 했습니다. (다음 슬라이드부터 자세히)"""
    add_text(slide1, 0.65, 3.7, 8.7, 1.6, problem, 9, False, RGBColor(0xE8, 0xE8, 0xE8))

    # ========== SLIDE 2: Traditional Weight Pipeline (시각화) ==========
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = LIGHT_BG
    bg2.line.fill.background()

    add_text(slide2, 0.5, 0.1, 9, 0.35, "기존 방식: RGB Weight Pipeline (상세 블록 다이어그램)", 15, True, DARK)

    # Old Pipeline Block Diagram using shapes
    # Box 1: RGB Input
    add_box(slide2, 0.5, 0.55, 1.8, 0.7, RGBColor(0xBB, 0xDE, 0xFB))
    add_text(slide2, 0.55, 0.6, 1.7, 0.6, "RGB 영상\n(작은 ROI)", 9, True, TEXT_DARK, align=PP_ALIGN.CENTER)

    add_arrow(slide2, 2.4, 0.75, 0.5, 0.3, RGBColor(0x90, 0xA4, 0xAE))

    # Box 2: Weight Learning
    add_box(slide2, 3.0, 0.55, 2.2, 0.7, RGBColor(0xFF, 0xCC, 0xBC))
    add_text(slide2, 3.05, 0.6, 2.1, 0.6, "RGB Weight 학습\n[w_r, w_g, w_b]\n(여러 번 실험)", 8, True, RED_LIGHT, align=PP_ALIGN.CENTER)

    add_arrow(slide2, 5.3, 0.75, 0.5, 0.3, RGBColor(0x90, 0xA4, 0xAE))

    # Box 3: Projection
    add_box(slide2, 5.9, 0.55, 1.6, 0.7, RGBColor(0xC8, 0xE6, 0xC9))
    add_text(slide2, 5.95, 0.6, 1.5, 0.6, "1D Pulse\n신호 생성", 9, True, TEXT_DARK, align=PP_ALIGN.CENTER)

    add_arrow(slide2, 7.6, 0.75, 0.5, 0.3, RGBColor(0x90, 0xA4, 0xAE))

    # Box 4: FFT + Peak
    add_box(slide2, 8.2, 0.55, 1.5, 0.7, RGBColor(0xFF, 0xF9, 0xC4))
    add_text(slide2, 8.25, 0.6, 1.4, 0.6, "FFT\n가장 큰\nPeak 선택", 8, True, TEXT_DARK, align=PP_ALIGN.CENTER)

    # Problem explanation
    add_box(slide2, 0.5, 1.4, 9, 1.8, RGBColor(0xFF, 0xEB, 0xEE))
    add_text(slide2, 0.65, 1.5, 8.7, 0.3, "이 방식의 치명적인 약점 (고심박 강아지에서)", 11, True, RGBColor(0xC0, 0x39, 0x2B))
    problem_detail = """1. Weight를 아무리 잘 찾아도, "한 번의 projection"으로 모든 정보를 1차원으로 압축한다.
   → 강아지 털 때문에 이미 손실이 크다.

2. FFT 후 "가장 큰 피크"를 고르는 것은, 스펙트럼 전체 모양을 무시하고 크기만 본다.
   → 고심박(200bpm)과 헐떡임 harmonic(100bpm)이 가까울 때 자주 실패.

3. Prior(이전 값 보정)는 "갑작스러운 큰 점프"를 막아주지만, 이미 잘못된 피크에 갇히면 벗어나기 어렵다.

실제 결과 (1/2/3 실험): 현실적 조건에서 최고 36 MAE. Video 3/7에서 특히 심각."""
    add_text(slide2, 0.65, 1.8, 8.7, 1.3, problem_detail, 9, False, TEXT_DARK)

    # Visual problem example
    add_text(slide2, 0.5, 3.35, 9, 0.3, "실제 실패 사례 (Video 3, target 210 bpm)", 11, True, TEAL)
    try:
        slide2.shapes.add_picture("reports/rppg_pet_keypoints/bvp_visualization/3_spectrum.png", 
                                  Inches(5.0), Inches(3.7), width=Inches(4.7))
    except:
        pass
    add_text(slide2, 0.5, 3.7, 4.4, 1.7, 
"""이 스펙트럼에서 기존 방법이 자주 126~148 bpm으로 예측한 이유:

- 파란색(100bpm 근처 artifact)이 세로로 더 높게 나온 window가 많음
- 컴퓨터는 "가장 높은 막대 = 심장"이라고 판단
- 결과: 실제 210 bpm인데 126 bpm으로 크게 틀림

이런 상황이 고심박 비디오에서 반복적으로 발생했습니다.""", 9, False, TEXT_DARK)

    # ========== SLIDE 3: New Spectrum Pipeline (상세 시각화) ==========
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = LIGHT_BG
    bg3.line.fill.background()

    add_text(slide3, 0.5, 0.1, 9, 0.35, "새로운 방식: Spectrum-Domain Pipeline (상세 블록 다이어그램)", 15, True, DARK)

    # New Pipeline Diagram
    # Step 1
    add_box(slide3, 0.3, 0.55, 1.5, 0.65, RGBColor(0xBB, 0xDE, 0xFB))
    add_text(slide3, 0.35, 0.58, 1.4, 0.6, "RGB 영상\n(ROI)", 8, True, TEXT_DARK, align=PP_ALIGN.CENTER)

    add_arrow(slide3, 1.9, 0.72, 0.4, 0.25, RGBColor(0x90, 0xA4, 0xAE))

    # Step 2: Fixed Views (no learning)
    add_box(slide3, 2.4, 0.55, 2.0, 0.65, RGBColor(0xDC, 0xED, 0xC8))
    add_text(slide3, 2.45, 0.58, 1.9, 0.6, "고정 Views\n(Green, G-R,\nPOS-like 등)\n→ 학습 없음!", 7, True, TEXT_DARK, align=PP_ALIGN.CENTER)

    add_arrow(slide3, 4.5, 0.72, 0.4, 0.25, RGBColor(0x90, 0xA4, 0xAE))

    # Step 3: Spectrum + Descriptors
    add_box(slide3, 5.0, 0.5, 2.3, 0.75, RGBColor(0xFF, 0xF9, 0xC4))
    add_text(slide3, 5.05, 0.52, 2.2, 0.7, "스펙트럼 변환\n+ Descriptors\n(peak 위치,\nHigh-HR vs Artifact\npower ratio)", 7, True, TEXT_DARK, align=PP_ALIGN.CENTER)

    add_arrow(slide3, 7.4, 0.72, 0.4, 0.25, RGBColor(0x90, 0xA4, 0xAE))

    # Step 4: Learned Model
    add_box(slide3, 7.9, 0.55, 1.8, 0.65, RGBColor(0xC8, 0xE6, 0xC9))
    add_text(slide3, 7.95, 0.58, 1.7, 0.6, "작은 모델\n(Ridge)\n→ BPM 직접\n예측", 8, True, SUCCESS, align=PP_ALIGN.CENTER)

    # Key difference highlight
    add_box(slide3, 0.5, 1.4, 9, 0.9, SOFT_GREEN)
    add_text(slide3, 0.65, 1.45, 8.7, 0.8, 
"""가장 큰 차이점: "RGB를 섞는 가중치(w_r, w_g, w_b)를 학습하지 않는다."
대신, 여러 고정된 방식으로 신호를 본 뒤, 그 신호들의 주파수 '모양' 전체를 feature로 사용합니다.
이게 강아지 털 환경과 고심박+헐떡임 상황에서 더 강력한 이유입니다.""", 9, False, TEXT_DARK)

    # Detailed explanation
    detail = """[Spectrum Descriptors가 실제로 계산하는 것들]

1. Binned Spectrum Power
   - 65~245 bpm 구간을 2.5 bpm 단위로 잘게 쪼갠다.
   - 각 구간에 에너지가 얼마나 몰려있는지 계산 (전체를 1로 정규화).

2. Peak Location
   - 가장 에너지가 강한 주파수가 어디인가? (200 근처인가, 100 근처인가?)

3. Power Ratio Features (가장 중요)
   - High-HR Cardiac Bands (165~185, 185~205, 205~225) vs Artifact Band (92~118)
   - "이번 그래프에서 고심박 구간이 헐떡임 구간보다 얼마나 더 두드러지는가?"
   - 이 비율이 높으면 → 심박수가 높을 가능성이 크다고 판단.

4. 추가 이진 힌트
   - Peak가 155 이상인가? (고심박일 가능성)

이 4가지를 조합한 feature vector를 작은 Ridge 모델에 넣으면,
"이런 모양의 스펙트럼이면 심박수는 대략 얼마"라고 직접 답을 줍니다.

기존처럼 "가장 큰 피크를 고르고 Prior로 보정"하는 게 아니라,
스펙트럼의 전체적인 '지문'을 보고 판단하는 방식입니다."""
    add_text(slide3, 0.5, 2.4, 9, 3.0, detail, 9, False, TEXT_DARK)

    # ========== SLIDE 4: Real Spectrum + Interpretation ==========
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    bg4 = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg4.fill.solid()
    bg4.fill.fore_color.rgb = LIGHT_BG
    bg4.line.fill.background()

    add_text(slide4, 0.5, 0.1, 9, 0.3, "실제 예시로 보는 차이 — Video 3 스펙트럼 상세 분석", 14, True, DARK)

    try:
        slide4.shapes.add_picture("reports/rppg_pet_keypoints/bvp_visualization/3_spectrum.png", 
                                  Inches(5.0), Inches(0.45), width=Inches(4.7))
    except:
        pass

    analysis = """왼쪽 그래프는 Video 3의 실제 한 구간을 주파수로 바꾼 것입니다.

[기존 Weight + Peak 방법이 실패하는 과정]
- 이 그래프에서 100bpm 근처(헐떡임 harmonic)에 매우 강한 에너지 피크가 있습니다.
- 200bpm 근처(실제 심박)에 에너지가 있지만, 상대적으로 약해 보일 때가 많습니다.
- 기존 방법은 "가장 높은 막대 = 심장"이라고 판단 → 126 bpm 또는 148 bpm으로 크게 틀림.

[Spectrum Descriptors가 성공하는 과정]
- Binned spectrum 전체를 보고 165~225 구간의 총 에너지를 계산.
- 92~118 artifact 구간과 비교 → 비율이 어느 정도인지 계산.
- 과거에 비슷한 비율을 보였던 데이터에서 실제 심박수가 어땠는지 학습.
- 결과: 218.9 bpm (실제 210 bpm에 매우 가까움).

이 차이가 1/2/3 실험에서 36 MAE → 18.2 MAE로 개선된 핵심 이유 중 하나입니다."""
    add_text(slide4, 0.5, 0.45, 4.4, 4.9, analysis, 9, False, TEXT_DARK)

    # ========== SLIDE 5: Kalman Tracker Visual ==========
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    bg5 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg5.fill.solid()
    bg5.fill.fore_color.rgb = LIGHT_BG
    bg5.line.fill.background()

    add_text(slide5, 0.5, 0.1, 9, 0.3, "Kalman State Tracker — 시간 흐름을 모델링하는 방법", 14, True, DARK)

    # Simple Kalman State Diagram using shapes
    # State 1: HR
    add_box(slide5, 1.5, 0.6, 2.5, 0.8, RGBColor(0xBB, 0xDE, 0xFB))
    add_text(slide5, 1.55, 0.65, 2.4, 0.7, "상태 1: 현재 HR\n(예: 185 bpm)", 9, True, TEXT_DARK, align=PP_ALIGN.CENTER)

    # State 2: Velocity
    add_box(slide5, 5.5, 0.6, 2.8, 0.8, RGBColor(0xC8, 0xE6, 0xC9))
    add_text(slide5, 5.55, 0.65, 2.7, 0.7, "상태 2: HR 변화 속도\n(Velocity, 예: +12 bpm/스텝)", 8, True, TEXT_DARK, align=PP_ALIGN.CENTER)

    add_text(slide5, 0.5, 1.5, 9, 0.3, "Tracker가 하는 일 (간단 비유)", 11, True, TEAL)
    tracker_text = """자동차 내비게이션이 "현재 속도 80km/h, 가속도 +30km/h"라고 알고 있으면,
갑자기 "280km/h"라는 이상한 측정값이 들어와도 "아, 이건 센서 오류일 가능성이 커"라고 판단하고
부드럽게 현실적인 값으로 보정해줍니다.

강아지 심박수도 마찬가지입니다.
- 심박수는 갑자기 50에서 200으로 뛰지 않습니다.
- 변화 속도(Velocity)는 시간이 지나면 점점 줄어듭니다 (안정화).

Kalman Tracker는 이 두 가지 사실을 수학적으로 모델링해서,
Spectrum이 가끔 이상한 값을 줘도 전체적으로 안정된 추정을 유지하게 만듭니다."""
    add_text(slide5, 0.5, 1.8, 9, 2.0, tracker_text, 9, False, TEXT_DARK)

    add_box(slide5, 0.5, 3.9, 9, 1.5, SOFT_GREEN)
    add_text(slide5, 0.65, 4.0, 8.7, 0.3, "Spectrum + Tracker를 함께 쓰는 이유", 11, True, SUCCESS)
    add_text(slide5, 0.65, 4.3, 8.7, 1.0, 
"""Spectrum Selector (1) → 순간순간 좋은 '측정값'과 대략적인 확신도를 제공
Kalman Tracker (3) → 그 측정값들을 시간 축에서 생리학적으로 그럴듯하게 연결

1+3 조합이 아직 18.2 (pure1)보다 낮게 나온 이유: Spectrum이 아직 자신의 '불확실도'를 정확히 알려주지 못해서 Tracker가 과도하게 부드럽게 만들었기 때문.
다음 단계: Spectrum 모델이 "이번 예측은 확신도가 낮다"는 신호를 잘 주도록 개선하면 1+3이 진짜 강력해집니다.""", 9, False, TEXT_DARK)

    # ========== SLIDE 6: Old vs New 완전 비교 + 결론 ==========
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    bg6 = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(5.625))
    bg6.fill.solid()
    bg6.fill.fore_color.rgb = DARK
    bg6.line.fill.background()

    add_text(slide6, 0.5, 0.15, 9, 0.35, "Old vs New: 알고리즘 패러다임 완전 비교", 16, True, WHITE)

    # Old box
    add_box(slide6, 0.5, 0.6, 4.3, 2.4, RGBColor(0x5D, 0x2E, 0x2E))
    add_text(slide6, 0.65, 0.7, 4.0, 0.3, "OLD: RGB Weight Paradigm", 11, True, RGBColor(0xFF, 0xCC, 0xCC))
    old_detail = """• 핵심 질문: "RGB를 어떻게 섞을까?"
• 방법: Weight 3개 학습 (v1, v2, high-HR 등)
• Peak 선택: FFT 후 가장 큰 값 + Prior 보정
• 강점: 간단하고 빠름
• 한계: 고심박+헐떡임에서 peak ambiguity 해결 불가
• 실험 결과: 현실적 조건 최고 36 MAE"""
    add_text(slide6, 0.65, 1.0, 4.0, 1.9, old_detail, 8, False, RGBColor(0xFF, 0xEB, 0xEE))

    # New box
    add_box(slide6, 5.2, 0.6, 4.3, 2.4, RGBColor(0x1E, 0x4D, 0x3A))
    add_text(slide6, 5.35, 0.7, 4.0, 0.3, "NEW: Spectrum + State Paradigm", 11, True, RGBColor(0xC8, 0xE6, 0xC9))
    new_detail = """• 핵심 질문: "스펙트럼 모양이 무엇을 말하는가?"
• 방법: 고정 views → Spectrum descriptors 학습
• Peak 선택: 모델이 전체 shape 보고 직접 예측
• 추가: Kalman으로 시간적 안정화
• 강점: artifact와 cardiac 경쟁 상황에서 강함
• 실험 결과: 18.2 MAE (Video 3에서 8.9)"""
    add_text(slide6, 5.35, 1.0, 4.0, 1.9, new_detail, 8, False, RGBColor(0xC8, 0xE6, 0xC9))

    # Final message
    add_box(slide6, 0.5, 3.15, 9, 2.3, RGBColor(0x0D, 0x1B, 0x2A))
    final = """[결론 — 발표에서 가장 강조할 부분]

우리는 1/2/3 실험을 통해 "RGB 가중치 3개를 더 잘 찾는 것"의 한계를 명확히 보았습니다.

그래서 질문을 바꿨습니다.
"강아지의 심박 신호는 어떤 고유한 모양을 가지고 있는가? 그 모양을 어떻게 읽어야 진짜를 알 수 있는가?"

그 질문에 답하기 위해:
- RGB를 1차원으로 압축하지 않고 주파수 영역의 모양 전체를 보기로 했습니다.
- 한 순간뿐만 아니라 시간에 따른 변화도 생리학적으로 모델링하기로 했습니다.

그 결과, 60개 데이터만으로도 이전에 50~70 bpm 오차가 나던 구간에서 8~9 bpm 수준까지 개선했습니다.

이것은 단순한 수치 개선이 아닙니다.
강아지 rPPG가 '인간 피부 모델의 변형'에서 '강아지 신호의 본질을 직접 읽는 모델'로 전환되는 첫 걸음입니다.

아직 완성된 길은 아니지만, 적어도 올바른 방향으로 첫걸음을 뗐습니다."""
    add_text(slide6, 0.65, 3.3, 8.7, 2.0, final, 9, False, WHITE)

    output = "reports/PET_RPPG_Iteration8_Visual_Detailed_Algorithm_Slides.pptx"
    prs.save(output)
    print(f"Generated: {output}")
    print("알고리즘을 그림(블록 다이어그램 + 실제 스펙트럼 이미지 + 플로우차트)과 함께 단계별로 매우 디테일하게 설명한 6장 슬라이드입니다.")

if __name__ == "__main__":
    main()